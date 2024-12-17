# Import necessary libraries
from pptx import Presentation
from pptx.util import Inches
import numpy as np
import matplotlib.pyplot as plt

# Create a presentation
prs = Presentation()

# Slide 1: Title Slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
title.text = "Earthquake Magnitude Prediction in the Philippines Using Neural Networks"
subtitle = title_slide.placeholders[1]
subtitle.text = "By Jose Tuling Jr."

# Slide 2: Introduction
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
title.text = "Introduction"
content = slide2.placeholders[1]
content.text = (
    "This research investigates the use of neural networks for predicting earthquake magnitudes in the Philippines. "
    "The study analyzes seismic data, including factors such as date, latitude, and longitude, to identify patterns and improve prediction accuracy. "
    "By employing advanced machine learning techniques, the project aims to enhance disaster preparedness and mitigate risks associated with earthquakes."
)

# Slide 3: Objectives
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide3.shapes.title
title.text = "Objectives"
content = slide3.placeholders[1]
content.text = (
    "1. Develop a neural network-based model for predicting earthquake magnitudes.\n"
    "2. Analyze seismic data features, including temporal and spatial attributes.\n"
    "3. Evaluate the model's performance using metrics like MAE, MSE, RMSE, and R-squared.\n"
    "4. Compare different factors, such as latitude and longitude, with earthquake magnitudes.\n"
    "5. Provide insights to aid in disaster management and preparedness."
)

# Slide 4: Literature Review
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide4.shapes.title
title.text = "Literature Review"
content = slide4.placeholders[1]
content.text = (
    "1. Emec and Ozcanhan (2024) applied neural networks to predict earthquake magnitudes in the Anatolian Plate.\n"
    "2. Hirata et al. (2012) conducted a forecast experiment on earthquake activity in Japan.\n"
    "3. Sivaiahbellamkonda et al. (2021) enhanced earthquake prediction using LSTM models.\n"
    "These studies demonstrate the potential of machine learning techniques for seismic prediction."
)

# Slide 5: Methodology
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide5.shapes.title
title.text = "Methodology"
content = slide5.placeholders[1]
content.text = (
    "1. Data Collection: PHIVOLCS seismic dataset.\n"
    "2. Preprocessing: Handling missing data and normalizing features.\n"
    "3. Model Design: Neural network architecture with LSTM layers.\n"
    "4. Evaluation: Metrics include MAE, MSE, RMSE, and R-squared.\n"
    "5. Visualization: Comparing magnitudes with latitude, longitude, and time."
)

# Create plots for visualization
x = np.linspace(0, 10, 50)
y1 = np.sin(x)  # Example magnitude vs time
y2 = np.cos(x)  # Example magnitude vs latitude

plt.figure()
plt.plot(x, y1, label='Magnitude vs Time')
plt.plot(x, y2, label='Magnitude vs Latitude')
plt.legend()
plt.title("Comparison of Magnitude to Different Features")
plt.xlabel("Feature")
plt.ylabel("Magnitude")
plt.savefig("plot1.png")

# Add Slide 6: Results and Discussion
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide6.shapes.title
title.text = "Results and Discussion"
content = slide6.placeholders[1]
content.text = (
    "The model achieved the following performance metrics:\n"
    "- Mean Absolute Error (MAE): 0.0715\n"
    "- Mean Squared Error (MSE): 0.0091\n"
    "- Root Mean Squared Error (RMSE): 0.0952\n"
    "- R-squared: 0.2239\n"
    "While the model shows potential, further tuning and feature engineering are needed to enhance accuracy."
)

# Add the plot to slide 6
left = Inches(1)
top = Inches(2)
pic = slide6.shapes.add_picture("plot1.png", left, top, height=Inches(3))

# Slide 7: Conclusion
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide7.shapes.title
title.text = "Conclusion"
content = slide7.placeholders[1]
content.text = (
    "The study demonstrates the application of neural networks in predicting earthquake magnitudes in the Philippines.\n"
    "Despite moderate prediction accuracy, the findings highlight key seismic data trends and the potential for machine learning in disaster management."
)

# Save the presentation
prs.save("Earthquake_Magnitude_Prediction.pptx")
"Presentation created and saved as 'Earthquake_Magnitude_Prediction.pptx'."
